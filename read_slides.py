import pptx
import pandas
import re


def extract_text(slide):
    text = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text.append(run.text)
    return ' '.join(text)


prs = pptx.Presentation('collection.pptx')

slides_df = pandas.DataFrame(columns=['raw content', 'raw roles', 'raw industries', 'raw projects', 'raw responsibilities', 'raw topics'], index=range(0, len(prs.slides)))

for i in range(len(prs.slides)):
    text = extract_text(prs.slides[i])
    slides_df.loc[i].at['raw content'] = text

for i in range(len(slides_df)):
    roles = re.findall(r'Role(.*?)Industry|Rolle(.*?)Industrie|Rolle(.*?)Branche', slides_df.iloc[i,0])
    slides_df.loc[i].at['raw roles'] = roles

for i in range(len(slides_df)):
    industries = re.findall(r'Industry(.*?)Project|Industrie(.*?)Projekt|Branche(.*?)Projekt', slides_df.iloc[i,0])
    slides_df.loc[i].at['raw industries'] = industries

for i in range(len(slides_df)):
    projects = re.findall(r'Project:(.*?)Responsible|Project :(.*?)Responsible|Project :(.*?)Responsibility|Projekt:(.*?)Verantwortlich|Projekt:(.*?)Aufgaben', slides_df.iloc[i,0])
    slides_df.loc[i].at['raw projects'] = projects

for i in range(len(slides_df)):
    responsibilities = re.findall(r'Verantwortlich(.*?)Rolle|Aufgaben(.*?)Rolle|Responsibility(.*?)Role|Responsible(.*?)Role|Responsible.*|Projekt:(.*?)Rolle', slides_df.iloc[i,0])
    slides_df.loc[i].at['raw responsibilities'] = responsibilities
